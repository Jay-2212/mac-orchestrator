import os
os.environ["HF_HUB_DISABLE_PROGRESS_BARS"] = "1"
os.environ["TOKENIZERS_PARALLELISM"] = "false"

import sqlite3
import json
import time
import uuid
import requests
import csv
import gc
import torch
from pathlib import Path
from sentence_transformers import SentenceTransformer
from concurrent.futures import ThreadPoolExecutor, as_completed

# Optional parsers (using try-except to allow partial runs if not all are installed)
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pptx
except ImportError:
    pptx = None

# --- CONFIGURATION ---
DIRECTORIES_TO_INDEX = [
    os.path.expanduser("~/Documents"),
    os.path.expanduser("~/Downloads"),
    os.path.expanduser("~/Desktop"),
    os.path.expanduser("~/OneDrive - Manipal Academy of Higher Education")
]

IGNORE_PATTERNS = {
    "node_modules", ".git", "library", ".venv", "venv", "__pycache__",
    "build", "dist", "target", "bin", "obj", "out", 
    ".gradle", ".idea", ".vscode", ".wrangler", "system-cache", 
    "processed", "cache", "buildinfo",
    "rss feed", "aeon-reading", "aeon reading", "varc", "asptor"
}
SUPPORTED_EXTENSIONS = {".txt", ".md", ".json", ".pdf", ".docx", ".xlsx", ".csv", ".pptx"}

# REPLACE THIS with your deployed Cloudflare Worker URL or localhost for testing
WORKER_URL = "https://mac-brain-worker.jb-brain.workers.dev"
INGEST_TOKEN = os.getenv("INGEST_TOKEN", "mac-brain-secret-key-123")

DB_PATH = os.path.expanduser("~/Documents/mac-orchestrator/sync_state.db")
CHUNK_SIZE = 500
OVERLAP = 50
BATCH_SIZE = 100
MAX_WORKERS = 8
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB — skip files larger than this to prevent OOM
MAX_FILES_PER_RUN = 30  # Limit to 30 files per process to prevent PyTorch/parsing RAM accumulation OOM

model = None

# --- PERSISTENT HTTP SESSION (reuses TCP+TLS connections across threads) ---
http_session = requests.Session()
http_session.headers.update({"Authorization": f"Bearer {INGEST_TOKEN}"})

# --- LOGGING HELPER ---
def log_msg(msg):
    """Print with timestamp for clear log ordering."""
    print(f"[{time.strftime('%H:%M:%S')}] {msg}")

# --- DATABASE SETUP ---
def init_db():
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS file_state 
                 (path TEXT PRIMARY KEY, last_modified REAL)''')
    conn.commit()
    return conn

# --- TEXT EXTRACTION & CHUNKING ---
def is_header_cell(val):
    if val is None:
        return False
    s = str(val).strip()
    if not s:
        return False
    try:
        float(s)
        return False
    except ValueError:
        pass
    return any(c.isalpha() for c in s)

def get_file_chunks(file_path):
    ext = Path(file_path).suffix.lower()
    
    # Skip oversized files to prevent OOM
    try:
        if os.path.getsize(file_path) > MAX_FILE_SIZE:
            log_msg(f"  Skipping oversized file ({os.path.getsize(file_path) // (1024*1024)}MB): {file_path}")
            return [] # Oversized: mark as processed since we deliberately skip it
    except OSError:
        return None
        
    chunks = []
    file_name = os.path.basename(file_path)
    
    try:
        if ext in {".txt", ".md"}:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            chunks = chunk_text_by_words(text, CHUNK_SIZE, OVERLAP)
            
        elif ext == ".json":
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                data = json.load(f)
                text = json.dumps(data, indent=2)
            chunks = chunk_text_by_words(text, CHUNK_SIZE, OVERLAP)
            
        elif ext == ".csv":
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                rows = list(reader)
            if rows:
                # Find first row that contains some textual content as header candidates
                header_row = None
                header_row_idx = 0
                for idx, r in enumerate(rows[:10]):
                    header_cells = [x for x in r if is_header_cell(x)]
                    if len(header_cells) >= 3:
                        header_row = r
                        header_row_idx = idx
                        break
                if header_row is None:
                    header_row = rows[0]
                    header_row_idx = 0
                
                headers = [str(h).strip() if h is not None else f"Col_{col_idx}" for col_idx, h in enumerate(header_row)]
                
                if len(rows) > 2000:
                    log_msg(f"  Large CSV file ({len(rows)} rows). Fallback to standard word-split chunking.")
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text = f.read()
                    chunks = chunk_text_by_words(text, CHUNK_SIZE, OVERLAP)
                else:
                    for row_idx, row in enumerate(rows[header_row_idx + 1:], start=header_row_idx + 1):
                        row_data = []
                        for col_idx, val in enumerate(row):
                            if val is not None and str(val).strip() != "":
                                header_name = headers[col_idx].strip() if col_idx < len(headers) else f"Col_{col_idx}"
                                row_data.append(f"{header_name}: {str(val).strip()}")
                        if row_data:
                            chunk_str = f"File: {file_name} | Row {row_idx + 1} -> " + " | ".join(row_data)
                            chunks.append(chunk_str)
                        
        elif ext == ".xlsx":
            if not openpyxl:
                log_msg(f"  Warning: openpyxl library is missing. Skipping: {file_path}")
                return None
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = list(sheet.iter_rows(values_only=True))
                if not rows:
                    continue
                
                # Detect header row (look at first 10 rows for descriptive headers containing letters)
                header_row = None
                header_row_idx = 0
                for idx, r in enumerate(rows[:10]):
                    header_cells = [x for x in r if is_header_cell(x)]
                    if len(header_cells) >= 3:
                        header_row = r
                        header_row_idx = idx
                        break
                
                if header_row is None:
                    header_row = rows[0]
                    header_row_idx = 0
                
                headers = [str(h).strip() if h is not None else f"Col_{col_idx}" for col_idx, h in enumerate(header_row)]
                
                if len(rows) > 2000:
                    log_msg(f"  Large worksheet {sheet_name} ({len(rows)} rows). Fallback to standard word-split chunking.")
                    sheet_text = ""
                    for r in rows:
                        row_data = [str(cell) for cell in r if cell is not None]
                        if row_data:
                            sheet_text += " | ".join(row_data) + "\n"
                    chunks.extend(chunk_text_by_words(sheet_text, CHUNK_SIZE, OVERLAP))
                else:
                    for idx, r in enumerate(rows[header_row_idx + 1:], start=header_row_idx + 1):
                        if not any(x is not None for x in r):
                            continue
                        row_data = []
                        for col_idx, val in enumerate(r):
                            if val is not None:
                                val_str = str(val).strip()
                                if val_str != "":
                                    header_name = headers[col_idx] if col_idx < len(headers) else f"Col_{col_idx}"
                                    header_name = header_name.replace("\n", " ").strip()
                                    row_data.append(f"{header_name}: {val_str}")
                        if row_data:
                            chunk_str = f"File: {file_name} | Sheet: {sheet_name}, Row {idx+1} -> " + " | ".join(row_data)
                            chunks.append(chunk_str)
                        
        elif ext == ".docx":
            if not docx:
                log_msg(f"  Warning: python-docx library is missing. Skipping: {file_path}")
                return None
            doc = docx.Document(file_path)
            elements_text = []
            
            for element in doc.element.body:
                if element.tag.endswith('p'):
                    p = docx.text.paragraph.Paragraph(element, doc)
                    if p.text.strip():
                        elements_text.append(p.text.strip())
                elif element.tag.endswith('tbl'):
                    t = docx.table.Table(element, doc)
                    t_rows = list(t.rows)
                    if t_rows:
                        headers = [cell.text.strip().replace("\n", " ") for cell in t_rows[0].cells]
                        for r_idx, r in enumerate(t_rows[1:], start=1):
                            row_cells = [cell.text.strip().replace("\n", " ") for cell in r.cells]
                            row_data = []
                            for col_idx, val in enumerate(row_cells):
                                if val:
                                    header_name = headers[col_idx] if col_idx < len(headers) else f"Col_{col_idx}"
                                    if not row_data or f"{header_name}: {val}" != row_data[-1]:
                                        row_data.append(f"{header_name}: {val}")
                            if row_data:
                                elements_text.append(f"Table Row -> " + " | ".join(row_data))
                                
            full_text = "\n".join(elements_text)
            chunks = chunk_text_by_words(full_text, CHUNK_SIZE, OVERLAP)
            
        elif ext == ".pptx":
            if not pptx:
                log_msg(f"  Warning: python-pptx library is missing. Skipping: {file_path}")
                return None
            prs = pptx.Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = f"File: {file_name} | Slide {slide_idx+1}"
                shapes_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        shapes_text.append(shape.text.strip().replace("\n", " "))
                if shapes_text:
                    slide_text += " -> " + " | ".join(shapes_text)
                    chunks.append(slide_text)
                    
        elif ext == ".pdf":
            if not PdfReader:
                log_msg(f"  Warning: pypdf library is missing. Skipping: {file_path}")
                return None
            reader = PdfReader(file_path)
            for page_idx, page in enumerate(reader.pages):
                extracted = page.extract_text()
                if extracted and extracted.strip():
                    page_text = f"File: {file_name} | Page {page_idx+1} -> {extracted.strip()}"
                    chunks.append(page_text)
                    
    except Exception as e:
        log_msg(f"  Error extracting {file_path}: {e}")
        return None
        
    return chunks

def chunk_text_by_words(text, chunk_size=CHUNK_SIZE, overlap=OVERLAP):
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
    return chunks

# --- MAIN INGESTION LOGIC ---
def run_indexer():
    conn = init_db()
    c = conn.cursor()
    
    files_to_process = []
    
    # 1. Discover Files
    log_msg("Scanning directories for new/modified files...")
    for directory in DIRECTORIES_TO_INDEX:
        if not os.path.exists(directory):
            continue
            
        for root, dirs, files in os.walk(directory):
            dirs[:] = [d for d in dirs if not any(pattern in d.lower() for pattern in IGNORE_PATTERNS)]
            
            for file in files:
                if file.startswith("~$"):
                    continue
                file_path = os.path.join(root, file)
                ext = Path(file_path).suffix.lower()
                
                if ext in SUPPORTED_EXTENSIONS:
                    try:
                        mtime = os.path.getmtime(file_path)
                        c.execute("SELECT last_modified FROM file_state WHERE path=?", (file_path,))
                        row = c.fetchone()
                        
                        if not row or row[0] < mtime:
                            files_to_process.append((file_path, mtime))
                    except Exception as e:
                        pass

    # Prioritize spreadsheets (.xlsx, .csv) at the beginning of the queue, with HOSPICE at the absolute top
    files_to_process.sort(key=lambda x: (
        "HOSPICE" not in os.path.basename(x[0]).upper(),
        Path(x[0]).suffix.lower() not in {".xlsx", ".csv"},
        Path(x[0]).suffix.lower(),
        x[0]
    ))

    total_found = len(files_to_process)
    log_msg(f"Found {total_found} new/modified files to index.")
    
    if MAX_FILES_PER_RUN and total_found > MAX_FILES_PER_RUN:
        log_msg(f"Limiting this run to the first {MAX_FILES_PER_RUN} files to manage memory footprint.")
        files_to_process = files_to_process[:MAX_FILES_PER_RUN]
    
    if not files_to_process:
        conn.close()
        log_msg("No changes detected. Exiting without loading model.")
        with open(os.path.expanduser("~/Documents/mac-orchestrator/indexer.log"), "a") as log_file:
            log_file.write(f"{time.strftime('%Y-%m-%d %H:%M:%S')} - SUCCESS: 0 files, model not loaded.\n")
        return
    
    # 2. Load model ONLY when files need processing (lazy loading)
    log_msg(f"Loading embedding model (CPU) for {len(files_to_process)} file(s)...")
    global model
    model = SentenceTransformer('BAAI/bge-base-en-v1.5', device='cpu')
    log_msg("Model loaded successfully.")
    
    failed_files_count = 0
    
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        for i, (file_path, mtime) in enumerate(files_to_process):
            log_msg(f"[{i+1}/{len(files_to_process)}] Processing: {file_path}")
            chunks = get_file_chunks(file_path)
            
            if chunks is None:
                # Extraction failed/skipped due to error or missing library.
                continue
                
            if not chunks:
                # Chunks are empty (e.g. empty file) — mark as processed so we don't retry.
                c.execute("INSERT OR REPLACE INTO file_state (path, last_modified) VALUES (?, ?)", (file_path, mtime))
                conn.commit()
                continue
            
            # Delete old chunks from Cloudflare DB first before uploading new ones to prevent orphan chunks
            delete_file_chunks(file_path)
            
            file_futures = []
            
            # Process chunks in sub-batches of BATCH_SIZE
            for chunk_start in range(0, len(chunks), BATCH_SIZE):
                chunk_subbatch = chunks[chunk_start : chunk_start + BATCH_SIZE]
                # Encode ONLY this small subbatch locally
                with torch.no_grad():
                    subbatch_vectors = model.encode(chunk_subbatch).tolist()
                
                upload_subbatch = []
                for idx, chunk in enumerate(chunk_subbatch):
                    chunk_idx = chunk_start + idx
                    chunk_id = f"{uuid.uuid5(uuid.NAMESPACE_URL, file_path)}-{chunk_idx}"
                    upload_subbatch.append({
                        "id": chunk_id,
                        "text": chunk,
                        "vector": subbatch_vectors[idx],
                        "path": file_path,
                        "modified": mtime
                    })
                
                # Submit the upload of this subbatch immediately
                file_futures.append(executor.submit(upload_batch, upload_subbatch))
                
                # Small cooling period to keep CPU low
                time.sleep(0.05)
            
            # Wait for all uploads of this specific file to complete
            file_failed = False
            for future in as_completed(file_futures):
                if not future.result():
                    file_failed = True
            
            if not file_failed:
                # Mark this specific file as completed in local SQLite and commit immediately!
                c.execute("INSERT OR REPLACE INTO file_state (path, last_modified) VALUES (?, ?)", (file_path, mtime))
                conn.commit()
                log_msg(f"  Successfully processed and saved progress for: {file_path}")
            else:
                log_msg(f"  WARNING: Some batches failed to upload for {file_path}. Will retry on next run.")
                failed_files_count += 1

            # Memory optimization: release large objects and clear PyTorch cache
            del chunks
            if 'upload_subbatch' in locals():
                del upload_subbatch
            if 'subbatch_vectors' in locals():
                del subbatch_vectors
            gc.collect()
            if torch.cuda.is_available():
                torch.cuda.empty_cache()
            if hasattr(torch, 'mps') and torch.mps.is_available():
                torch.mps.empty_cache()
                
    conn.close()
    status = f"{len(files_to_process)} files, {failed_files_count} failed files"
    log_msg(f"[{'SUCCESS' if failed_files_count == 0 else 'PARTIAL'}] Indexing complete: {status}")
    with open(os.path.expanduser("~/Documents/mac-orchestrator/indexer.log"), "a") as log_file:
        log_file.write(f"{time.strftime('%Y-%m-%d %H:%M:%S')} - {'SUCCESS' if failed_files_count == 0 else 'PARTIAL'}: {status}\n")

def delete_file_chunks(file_path):
    """Delete all existing chunks for a file from D1 and Vectorize before uploading new ones."""
    log_msg(f"  Purging existing chunks in Cloudflare for: {file_path}")
    try:
        resp = http_session.post(
            f"{WORKER_URL}/delete-file",
            json={"path": file_path},
            timeout=30
        )
        if resp.status_code == 200:
            res_json = resp.json()
            log_msg(f"  Successfully deleted {res_json.get('deletedCount', 0)} old chunks.")
            return True
        else:
            log_msg(f"  Failed to delete old chunks: {resp.status_code} - {resp.text}")
    except Exception as e:
        log_msg(f"  Network error deleting old chunks: {e}")
    return False

def upload_batch(batch):
    """Upload a batch of chunks to Cloudflare with retries. Returns True on success, False on failure."""
    log_msg(f"  Uploading batch of {len(batch)} chunks...")
    max_retries = 3
    delay = 1.0
    for attempt in range(1, max_retries + 1):
        try:
            resp = http_session.post(
                f"{WORKER_URL}/ingest",
                json=batch,
                timeout=30
            )
            if resp.status_code == 200:
                log_msg(f"  Batch uploaded successfully on attempt {attempt}.")
                return True
            else:
                log_msg(f"  Failed to upload batch (attempt {attempt}/{max_retries}): {resp.status_code} - {resp.text}")
        except Exception as e:
            log_msg(f"  Network error during upload (attempt {attempt}/{max_retries}): {e}")
        
        if attempt < max_retries:
            time.sleep(delay)
            delay *= 2
            
    return False

if __name__ == "__main__":
    run_indexer()
